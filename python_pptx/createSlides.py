#!/usr/bin/python
from pptx import Presentation
import sys
import random

from pptx.enum.shapes import MSO_SHAPE
prs = Presentation()

def title_slide(presentation,title,subtitle = ""):
	slide = prs.slides.add_slide(prs.slide_layouts[0])
	slide.shapes.title.text = title
	slide.placeholders[1].text = subtitle

def add_bullet_slide(presentation,title,bullets_text):
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	slide.shapes.title.text = title
	tf = slide.shapes.placeholders[1].textframe
	paragraph = tf
	for i in bullets_text:
		paragraph.text = i
		paragraph = tf.add_paragraph()

#read the source questions 
source = open(sys.argv[1],"r")
lines = source.readlines()
title_slide(prs,lines[0],lines[1])
lines = lines[2:len(lines)]
#remove the line breaks
for i in range(len(lines)):
	lines[i] = lines[i].replace("\n","")
#create questions from lines
questions = []
for l in range(len(lines)/5):
	questions.append((lines[l*5],lines[l*5+1:l*5+5]))
#randomize questions
random.shuffle(questions)
for q in questions:
	right = q[1][0]
	#randomize answers
	random.shuffle(q[1])
	add_bullet_slide(prs,q[0],q[1])

#save presentation
prs.save(sys.argv[2])
